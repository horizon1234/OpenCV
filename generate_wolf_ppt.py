#!/usr/bin/env python3
"""
生成"狼来了"故事 PPT（幼儿园版）

使用 Pillow 生成每页插图，再用 python-pptx 组装为 ~20 页 PPT。
运行: python generate_wolf_ppt.py
输出: 狼来了.pptx
"""

import math
import os
import random
from io import BytesIO

from lxml import etree
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# 常量
# ---------------------------------------------------------------------------
IMG_W, IMG_H = 960, 540  # 16:9 插图尺寸
OUTPUT_FILE = "狼来了.pptx"

# 颜色常量
SKY_BLUE = (135, 206, 235)
GRASS_GREEN = (34, 139, 34)
LIGHT_GREEN = (144, 238, 144)
SUN_YELLOW = (255, 223, 0)
SKIN = (255, 224, 189)
BROWN = (139, 69, 19)
DARK_BROWN = (101, 67, 33)
WHITE = (255, 255, 255)
BLACK = (0, 0, 0)
RED = (220, 20, 60)
GRAY = (128, 128, 128)
DARK_GRAY = (80, 80, 80)
ORANGE = (255, 165, 0)
NIGHT_BLUE = (25, 25, 112)
DARK_GREEN = (0, 100, 0)
PINK = (255, 182, 193)
LIGHT_YELLOW = (255, 255, 224)
WOLF_GRAY = (105, 105, 105)


# ---------------------------------------------------------------------------
# 绘图辅助
# ---------------------------------------------------------------------------

def _sky_and_grass(draw: ImageDraw.Draw, sky_color=SKY_BLUE):
    """画天空和草地背景。"""
    draw.rectangle([0, 0, IMG_W, IMG_H * 2 // 3], fill=sky_color)
    draw.rectangle([0, IMG_H * 2 // 3, IMG_W, IMG_H], fill=GRASS_GREEN)


def _sun(draw: ImageDraw.Draw, x=860, y=60, r=40):
    """画太阳。"""
    draw.ellipse([x - r, y - r, x + r, y + r], fill=SUN_YELLOW)
    for angle in range(0, 360, 30):
        x1 = x + int((r + 15) * math.cos(math.radians(angle)))
        y1 = y + int((r + 15) * math.sin(math.radians(angle)))
        x2 = x + int((r + 30) * math.cos(math.radians(angle)))
        y2 = y + int((r + 30) * math.sin(math.radians(angle)))
        draw.line([x1, y1, x2, y2], fill=SUN_YELLOW, width=3)


def _cloud(draw: ImageDraw.Draw, x, y):
    """画一朵云。"""
    for dx, dy, r in [(0, 0, 25), (20, -10, 20), (-20, -8, 18), (10, 5, 22)]:
        draw.ellipse([x + dx - r, y + dy - r, x + dx + r, y + dy + r], fill=WHITE)


def _tree(draw: ImageDraw.Draw, x, y):
    """画一棵简单的树。"""
    draw.rectangle([x - 8, y, x + 8, y + 50], fill=BROWN)
    draw.ellipse([x - 35, y - 40, x + 35, y + 20], fill=DARK_GREEN)


def _sheep(draw: ImageDraw.Draw, x, y, size=1.0):
    """画一只简笔小绵羊。"""
    s = size
    # 身体 (白色椭圆)
    draw.ellipse([x - int(25*s), y - int(15*s), x + int(25*s), y + int(15*s)], fill=WHITE, outline=BLACK, width=2)
    # 头
    draw.ellipse([x + int(20*s), y - int(20*s), x + int(38*s), y - int(2*s)], fill=WHITE, outline=BLACK, width=2)
    # 眼睛
    draw.ellipse([x + int(28*s), y - int(15*s), x + int(32*s), y - int(11*s)], fill=BLACK)
    # 腿
    for lx in [-15, -5, 5, 15]:
        draw.line([x + int(lx*s), y + int(15*s), x + int(lx*s), y + int(30*s)], fill=BLACK, width=2)


def _boy(draw: ImageDraw.Draw, x, y, mouth_open=False, arms_up=False):
    """画放羊小孩。"""
    # 身体
    draw.rectangle([x - 12, y, x + 12, y + 40], fill=RED)
    # 头
    draw.ellipse([x - 15, y - 30, x + 15, y], fill=SKIN)
    # 眼睛
    draw.ellipse([x - 8, y - 20, x - 4, y - 14], fill=BLACK)
    draw.ellipse([x + 4, y - 20, x + 8, y - 14], fill=BLACK)
    # 嘴巴
    if mouth_open:
        draw.ellipse([x - 5, y - 10, x + 5, y - 2], fill=BLACK)
    else:
        draw.arc([x - 5, y - 12, x + 5, y - 4], start=0, end=180, fill=BLACK, width=2)
    # 头发
    draw.arc([x - 15, y - 35, x + 15, y - 10], start=180, end=360, fill=DARK_BROWN, width=4)
    # 手臂
    if arms_up:
        draw.line([x - 12, y + 5, x - 30, y - 15], fill=SKIN, width=4)
        draw.line([x + 12, y + 5, x + 30, y - 15], fill=SKIN, width=4)
    else:
        draw.line([x - 12, y + 5, x - 25, y + 25], fill=SKIN, width=4)
        draw.line([x + 12, y + 5, x + 25, y + 25], fill=SKIN, width=4)
    # 腿
    draw.line([x - 6, y + 40, x - 10, y + 60], fill=BROWN, width=4)
    draw.line([x + 6, y + 40, x + 10, y + 60], fill=BROWN, width=4)


def _villager(draw: ImageDraw.Draw, x, y, color=(0, 0, 200)):
    """画一个村民。"""
    # 身体
    draw.rectangle([x - 10, y, x + 10, y + 35], fill=color)
    # 头
    draw.ellipse([x - 12, y - 25, x + 12, y], fill=SKIN)
    # 眼睛
    draw.ellipse([x - 6, y - 17, x - 2, y - 12], fill=BLACK)
    draw.ellipse([x + 2, y - 17, x + 6, y - 12], fill=BLACK)
    # 嘴
    draw.arc([x - 4, y - 10, x + 4, y - 4], start=0, end=180, fill=BLACK, width=2)
    # 腿
    draw.line([x - 5, y + 35, x - 8, y + 52], fill=BROWN, width=3)
    draw.line([x + 5, y + 35, x + 8, y + 52], fill=BROWN, width=3)


def _wolf(draw: ImageDraw.Draw, x, y, size=1.0):
    """画一只大灰狼。"""
    s = size
    # 身体
    draw.ellipse([x - int(35*s), y - int(18*s), x + int(35*s), y + int(18*s)], fill=WOLF_GRAY, outline=DARK_GRAY, width=2)
    # 头
    draw.ellipse([x + int(25*s), y - int(30*s), x + int(55*s), y - int(5*s)], fill=WOLF_GRAY, outline=DARK_GRAY, width=2)
    # 耳朵
    draw.polygon([
        (x + int(32*s), y - int(30*s)),
        (x + int(28*s), y - int(45*s)),
        (x + int(38*s), y - int(30*s)),
    ], fill=WOLF_GRAY, outline=DARK_GRAY)
    draw.polygon([
        (x + int(42*s), y - int(30*s)),
        (x + int(45*s), y - int(48*s)),
        (x + int(50*s), y - int(30*s)),
    ], fill=WOLF_GRAY, outline=DARK_GRAY)
    # 眼睛 (红色凶恶)
    draw.ellipse([x + int(35*s), y - int(24*s), x + int(42*s), y - int(17*s)], fill=RED)
    draw.ellipse([x + int(38*s), y - int(22*s), x + int(40*s), y - int(19*s)], fill=BLACK)
    # 鼻子
    draw.ellipse([x + int(50*s), y - int(20*s), x + int(55*s), y - int(15*s)], fill=BLACK)
    # 嘴巴 (张嘴)
    draw.arc([x + int(40*s), y - int(15*s), x + int(55*s), y - int(5*s)], start=0, end=180, fill=BLACK, width=2)
    # 尾巴
    draw.arc([x - int(35*s), y - int(25*s), x - int(15*s), y + int(5*s)], start=180, end=320, fill=WOLF_GRAY, width=int(6*s))
    # 腿
    for lx in [-20, -8, 8, 20]:
        draw.line([x + int(lx*s), y + int(18*s), x + int(lx*s), y + int(35*s)], fill=DARK_GRAY, width=int(3*s))


def _mountain(draw: ImageDraw.Draw):
    """画远处的山。"""
    draw.polygon([(0, 360), (200, 200), (400, 360)], fill=(100, 160, 100))
    draw.polygon([(300, 360), (550, 180), (750, 360)], fill=(80, 140, 80))


def _fence(draw: ImageDraw.Draw, x, y, count=5):
    """画栅栏。"""
    for i in range(count):
        fx = x + i * 20
        draw.rectangle([fx, y - 30, fx + 5, y + 10], fill=BROWN, outline=DARK_BROWN)
    draw.rectangle([x - 5, y - 20, x + count * 20, y - 15], fill=BROWN, outline=DARK_BROWN)
    draw.rectangle([x - 5, y - 5, x + count * 20, y], fill=BROWN, outline=DARK_BROWN)


def _stars(draw: ImageDraw.Draw, count=20):
    """画星星（夜景用）。"""
    rng = random.Random(42)
    for _ in range(count):
        sx = rng.randint(20, IMG_W - 20)
        sy = rng.randint(20, IMG_H * 2 // 3 - 20)
        draw.ellipse([sx - 2, sy - 2, sx + 2, sy + 2], fill=SUN_YELLOW)


def _moon(draw: ImageDraw.Draw, x=100, y=60):
    """画月亮。"""
    draw.ellipse([x - 30, y - 30, x + 30, y + 30], fill=LIGHT_YELLOW)
    draw.ellipse([x - 15, y - 35, x + 25, y + 25], fill=NIGHT_BLUE)


def _angry_villager(draw: ImageDraw.Draw, x, y, color=(0, 0, 200)):
    """画一个生气的村民。"""
    # 身体
    draw.rectangle([x - 10, y, x + 10, y + 35], fill=color)
    # 头
    draw.ellipse([x - 12, y - 25, x + 12, y], fill=SKIN)
    # 愤怒的眉毛
    draw.line([x - 8, y - 20, x - 2, y - 17], fill=BLACK, width=2)
    draw.line([x + 2, y - 17, x + 8, y - 20], fill=BLACK, width=2)
    # 眼睛
    draw.ellipse([x - 7, y - 16, x - 3, y - 12], fill=BLACK)
    draw.ellipse([x + 3, y - 16, x + 7, y - 12], fill=BLACK)
    # 嘴 (不高兴)
    draw.arc([x - 4, y - 8, x + 4, y - 2], start=180, end=360, fill=BLACK, width=2)
    # 腿
    draw.line([x - 5, y + 35, x - 8, y + 52], fill=BROWN, width=3)
    draw.line([x + 5, y + 35, x + 8, y + 52], fill=BROWN, width=3)


def _sad_boy(draw: ImageDraw.Draw, x, y):
    """画哭泣的小男孩。"""
    # 身体
    draw.rectangle([x - 12, y, x + 12, y + 40], fill=RED)
    # 头
    draw.ellipse([x - 15, y - 30, x + 15, y], fill=SKIN)
    # 眼睛 (闭着, 哭)
    draw.arc([x - 9, y - 20, x - 3, y - 14], start=0, end=180, fill=BLACK, width=2)
    draw.arc([x + 3, y - 20, x + 9, y - 14], start=0, end=180, fill=BLACK, width=2)
    # 泪水
    draw.ellipse([x - 10, y - 13, x - 7, y - 8], fill=(0, 150, 255))
    draw.ellipse([x + 7, y - 13, x + 10, y - 8], fill=(0, 150, 255))
    # 嘴 (难过)
    draw.arc([x - 5, y - 8, x + 5, y - 2], start=180, end=360, fill=BLACK, width=2)
    # 头发
    draw.arc([x - 15, y - 35, x + 15, y - 10], start=180, end=360, fill=DARK_BROWN, width=4)
    # 手臂 (垂下)
    draw.line([x - 12, y + 5, x - 25, y + 25], fill=SKIN, width=4)
    draw.line([x + 12, y + 5, x + 25, y + 25], fill=SKIN, width=4)
    # 腿
    draw.line([x - 6, y + 40, x - 10, y + 60], fill=BROWN, width=4)
    draw.line([x + 6, y + 40, x + 10, y + 60], fill=BROWN, width=4)


# ---------------------------------------------------------------------------
# 每一页的场景绘制函数
# ---------------------------------------------------------------------------

def draw_cover(img: Image.Image):
    """封面。"""
    draw = ImageDraw.Draw(img)
    _sky_and_grass(draw)
    _sun(draw)
    _cloud(draw, 150, 80)
    _cloud(draw, 450, 50)
    _mountain(draw)
    _boy(draw, 350, 300)
    _sheep(draw, 500, 340)
    _sheep(draw, 600, 350)
    _wolf(draw, 150, 380, size=0.9)


def draw_scene_village(img: Image.Image):
    """场景：美丽的小村庄。"""
    draw = ImageDraw.Draw(img)
    _sky_and_grass(draw)
    _sun(draw)
    _cloud(draw, 200, 70)
    _mountain(draw)
    # 房子
    for hx in [150, 400, 650]:
        draw.rectangle([hx, 280, hx + 80, 360], fill=ORANGE)
        draw.polygon([(hx - 10, 280), (hx + 40, 240), (hx + 90, 280)], fill=RED)
        draw.rectangle([hx + 30, 320, hx + 50, 360], fill=BROWN)
    # 树
    _tree(draw, 80, 300)
    _tree(draw, 800, 310)


def draw_scene_boy_intro(img: Image.Image):
    """场景：介绍放羊娃。"""
    draw = ImageDraw.Draw(img)
    _sky_and_grass(draw)
    _sun(draw)
    _cloud(draw, 300, 60)
    _boy(draw, 250, 310)
    _sheep(draw, 400, 350)
    _sheep(draw, 500, 340)
    _sheep(draw, 600, 355)
    _fence(draw, 350, 400, 8)


def draw_scene_bored(img: Image.Image):
    """场景：放羊娃觉得无聊。"""
    draw = ImageDraw.Draw(img)
    _sky_and_grass(draw)
    _sun(draw)
    _cloud(draw, 500, 50)
    _boy(draw, 300, 310)
    # 无聊的气泡
    draw.ellipse([320, 250, 420, 290], fill=WHITE, outline=BLACK, width=2)
    draw.text((340, 260), "...", fill=BLACK)
    _sheep(draw, 550, 350)
    _sheep(draw, 650, 345)
    _sheep(draw, 480, 360)


def draw_scene_idea(img: Image.Image):
    """场景：放羊娃想到坏主意。"""
    draw = ImageDraw.Draw(img)
    _sky_and_grass(draw)
    _sun(draw)
    _boy(draw, 300, 310)
    # 灯泡气泡
    draw.ellipse([310, 240, 390, 290], fill=LIGHT_YELLOW, outline=BLACK, width=2)
    draw.ellipse([340, 250, 360, 275], fill=SUN_YELLOW)
    draw.line([350, 275, 350, 285], fill=DARK_GRAY, width=2)
    _sheep(draw, 550, 350)
    _sheep(draw, 650, 345)


def draw_scene_first_cry(img: Image.Image):
    """场景：第一次喊"狼来了！"。"""
    draw = ImageDraw.Draw(img)
    _sky_and_grass(draw)
    _sun(draw)
    _cloud(draw, 200, 60)
    _boy(draw, 300, 310, mouth_open=True, arms_up=True)
    # 喊话气泡
    draw.ellipse([330, 230, 500, 280], fill=WHITE, outline=RED, width=3)
    draw.text((355, 245), "狼来了!", fill=RED)
    _sheep(draw, 550, 350)
    _sheep(draw, 650, 345)


def draw_scene_villagers_run(img: Image.Image):
    """场景：村民们赶来。"""
    draw = ImageDraw.Draw(img)
    _sky_and_grass(draw)
    _sun(draw)
    _villager(draw, 600, 320, (0, 0, 200))
    _villager(draw, 680, 315, (200, 0, 0))
    _villager(draw, 760, 325, (0, 150, 0))
    # 拿工具
    draw.line([615, 320, 635, 290], fill=BROWN, width=3)  # 棍子
    draw.line([695, 315, 715, 285], fill=GRAY, width=3)   # 锄头
    _boy(draw, 300, 310, mouth_open=True, arms_up=True)
    _sheep(draw, 450, 350)


def draw_scene_no_wolf(img: Image.Image):
    """场景：没有狼，放羊娃大笑。"""
    draw = ImageDraw.Draw(img)
    _sky_and_grass(draw)
    _sun(draw)
    _cloud(draw, 400, 50)
    _boy(draw, 300, 310, mouth_open=True)
    # 笑的符号
    draw.text((280, 260), "哈哈哈", fill=RED)
    _villager(draw, 500, 320, (0, 0, 200))
    _villager(draw, 580, 315, (200, 0, 0))
    # 问号
    draw.text((505, 280), "?", fill=BLACK)
    draw.text((585, 275), "?", fill=BLACK)
    _sheep(draw, 700, 350)


def draw_scene_villagers_leave(img: Image.Image):
    """场景：村民们生气离开。"""
    draw = ImageDraw.Draw(img)
    _sky_and_grass(draw)
    _sun(draw)
    _angry_villager(draw, 600, 320, (0, 0, 200))
    _angry_villager(draw, 700, 315, (200, 0, 0))
    _angry_villager(draw, 800, 325, (0, 150, 0))
    _boy(draw, 250, 310, mouth_open=True)
    draw.text((230, 260), "哈哈", fill=RED)
    _sheep(draw, 400, 350)


def draw_scene_second_cry(img: Image.Image):
    """场景：第二次喊"狼来了！"。"""
    draw = ImageDraw.Draw(img)
    _sky_and_grass(draw)
    _sun(draw, y=80)
    _cloud(draw, 700, 60)
    _boy(draw, 300, 310, mouth_open=True, arms_up=True)
    draw.ellipse([330, 225, 520, 275], fill=WHITE, outline=RED, width=3)
    draw.text((350, 238), "狼来了!!", fill=RED)
    _sheep(draw, 550, 350)
    _sheep(draw, 650, 345)


def draw_scene_villagers_run2(img: Image.Image):
    """场景：村民们又赶来了。"""
    draw = ImageDraw.Draw(img)
    _sky_and_grass(draw)
    _sun(draw)
    _villager(draw, 600, 320, (0, 0, 200))
    _villager(draw, 680, 315, (200, 0, 0))
    _villager(draw, 760, 325, (0, 150, 0))
    _boy(draw, 300, 310, mouth_open=True, arms_up=True)
    _sheep(draw, 450, 350)
    _cloud(draw, 200, 50)


def draw_scene_no_wolf2(img: Image.Image):
    """场景：又没有狼。"""
    draw = ImageDraw.Draw(img)
    _sky_and_grass(draw)
    _sun(draw)
    _boy(draw, 300, 310, mouth_open=True)
    draw.text((275, 260), "哈哈哈哈", fill=RED)
    _angry_villager(draw, 520, 320, (0, 0, 200))
    _angry_villager(draw, 620, 315, (200, 0, 0))
    _sheep(draw, 700, 350)


def draw_scene_angry_leave2(img: Image.Image):
    """场景：村民非常生气。"""
    draw = ImageDraw.Draw(img)
    _sky_and_grass(draw)
    _sun(draw)
    _angry_villager(draw, 550, 320, (0, 0, 200))
    _angry_villager(draw, 650, 315, (200, 0, 0))
    _angry_villager(draw, 750, 320, (0, 150, 0))
    # 生气的符号
    draw.text((555, 280), "╬", fill=RED)
    draw.text((660, 275), "╬", fill=RED)
    _boy(draw, 250, 310)
    _sheep(draw, 400, 350)


def draw_scene_wolf_comes(img: Image.Image):
    """场景：狼真的来了！"""
    draw = ImageDraw.Draw(img)
    _sky_and_grass(draw, sky_color=(200, 150, 150))
    _cloud(draw, 300, 60)
    _wolf(draw, 700, 350, size=1.2)
    _sheep(draw, 450, 350)
    _sheep(draw, 350, 360)
    _boy(draw, 200, 310)
    # 惊恐
    draw.text((170, 260), "!!!", fill=RED)


def draw_scene_boy_cries_help(img: Image.Image):
    """场景：放羊娃拼命喊救命。"""
    draw = ImageDraw.Draw(img)
    _sky_and_grass(draw, sky_color=(200, 150, 150))
    _boy(draw, 250, 310, mouth_open=True, arms_up=True)
    draw.ellipse([270, 225, 460, 275], fill=WHITE, outline=RED, width=3)
    draw.text((285, 237), "救命啊!", fill=RED)
    _wolf(draw, 650, 350, size=1.2)
    _sheep(draw, 450, 360)


def draw_scene_nobody_comes(img: Image.Image):
    """场景：没有人来帮忙。"""
    draw = ImageDraw.Draw(img)
    _sky_and_grass(draw, sky_color=(180, 140, 140))
    _sad_boy(draw, 200, 310)
    _wolf(draw, 600, 340, size=1.2)
    _sheep(draw, 400, 360)
    # 远处的村子没人出来
    for hx in [700, 800]:
        draw.rectangle([hx, 300, hx + 50, 360], fill=ORANGE)
        draw.polygon([(hx - 5, 300), (hx + 25, 270), (hx + 55, 300)], fill=RED)


def draw_scene_wolf_attack(img: Image.Image):
    """场景：狼抓走了羊。"""
    draw = ImageDraw.Draw(img)
    _sky_and_grass(draw, sky_color=(180, 140, 140))
    _wolf(draw, 500, 340, size=1.3)
    # 被叼走的羊
    draw.ellipse([540, 310, 580, 335], fill=WHITE, outline=BLACK, width=2)
    _sad_boy(draw, 200, 310)


def draw_scene_boy_sad(img: Image.Image):
    """场景：放羊娃非常后悔。"""
    draw = ImageDraw.Draw(img)
    _sky_and_grass(draw, sky_color=(170, 190, 220))
    _sad_boy(draw, 480, 310)
    # 空旷的山坡，没有羊了
    _fence(draw, 200, 400, 8)
    _tree(draw, 100, 300)
    _cloud(draw, 600, 60)


def draw_scene_lesson(img: Image.Image):
    """场景：道理 —— 诚实很重要。"""
    draw = ImageDraw.Draw(img)
    _sky_and_grass(draw)
    _sun(draw, x=480, y=80, r=50)
    _cloud(draw, 150, 50)
    _cloud(draw, 750, 70)
    # 小男孩和村民和好
    _boy(draw, 350, 320)
    _villager(draw, 550, 320, (0, 0, 200))
    _villager(draw, 650, 315, (200, 0, 0))
    # 爱心
    draw.polygon([(450, 290), (460, 280), (470, 290), (460, 305)], fill=RED)


def draw_ending(img: Image.Image):
    """结尾页。"""
    draw = ImageDraw.Draw(img)
    _sky_and_grass(draw)
    _sun(draw, x=480, y=100, r=55)
    _cloud(draw, 100, 60)
    _cloud(draw, 800, 80)
    _tree(draw, 100, 300)
    _tree(draw, 850, 310)
    _sheep(draw, 400, 370)
    _sheep(draw, 500, 375)
    _sheep(draw, 600, 365)
    _boy(draw, 300, 330)
    _fence(draw, 350, 420, 10)


# ---------------------------------------------------------------------------
# 幻灯片内容定义
# ---------------------------------------------------------------------------

SLIDES = [
    {
        "draw": draw_cover,
        "title": "狼 来 了",
        "body": "——中国经典寓言故事——\n幼儿园绘本",
        "title_size": 44,
        "body_size": 24,
    },
    {
        "draw": draw_scene_village,
        "title": "",
        "body": "从前，在一个美丽的小村庄旁边，\n有一座绿绿的大山。",
        "body_size": 26,
    },
    {
        "draw": draw_scene_boy_intro,
        "title": "",
        "body": '村子里有一个小男孩，\n大家叫他"放羊娃"。\n他每天都要到山上去放羊。🐑',
        "body_size": 26,
    },
    {
        "draw": draw_scene_bored,
        "title": "",
        "body": "山上只有他一个人，\n没有小朋友跟他玩，\n他觉得好无聊啊……",
        "body_size": 26,
    },
    {
        "draw": draw_scene_idea,
        "title": "",
        "body": '忽然，他想到了一个\n"好玩"的主意……💡',
        "body_size": 26,
    },
    {
        "draw": draw_scene_first_cry,
        "title": "",
        "body": '他站起来大喊：\n"狼来了！狼来了！\n快来救我呀！"🗣️',
        "body_size": 26,
    },
    {
        "draw": draw_scene_villagers_run,
        "title": "",
        "body": "山下的村民们听到了，\n赶紧拿起锄头和棍子，\n跑上山来帮他！",
        "body_size": 26,
    },
    {
        "draw": draw_scene_no_wolf,
        "title": "",
        "body": '可是，山上根本没有狼！\n放羊娃哈哈大笑：\n"骗你们的！哈哈哈！"',
        "body_size": 26,
    },
    {
        "draw": draw_scene_villagers_leave,
        "title": "",
        "body": "村民们很生气，\n摇摇头就走了。😤",
        "body_size": 26,
    },
    {
        "draw": draw_scene_second_cry,
        "title": "",
        "body": '过了几天，放羊娃又觉得无聊了。\n他又大喊起来：\n"狼来了！狼来了！！"',
        "body_size": 26,
    },
    {
        "draw": draw_scene_villagers_run2,
        "title": "",
        "body": "村民们又急急忙忙跑上了山。\n大家气喘吁吁的……",
        "body_size": 26,
    },
    {
        "draw": draw_scene_no_wolf2,
        "title": "",
        "body": '山上还是没有狼！\n放羊娃又笑了：\n"又上当了！哈哈哈哈！"',
        "body_size": 26,
    },
    {
        "draw": draw_scene_angry_leave2,
        "title": "",
        "body": '这一次，村民们非常非常生气。\n他们说："我们再也不相信你了！"',
        "body_size": 26,
    },
    {
        "draw": draw_scene_wolf_comes,
        "title": "",
        "body": "又过了几天……\n这一次，一只大灰狼\n真的来了！！！🐺",
        "body_size": 26,
    },
    {
        "draw": draw_scene_boy_cries_help,
        "title": "",
        "body": '放羊娃吓坏了！\n他拼命喊：\n"狼来了！狼真的来了！\n快来救我啊！"',
        "body_size": 26,
    },
    {
        "draw": draw_scene_nobody_comes,
        "title": "",
        "body": "可是这一次……\n没有一个人来帮他。😢\n大家都以为他又在骗人。",
        "body_size": 26,
    },
    {
        "draw": draw_scene_wolf_attack,
        "title": "",
        "body": "大灰狼抓走了好几只小羊。\n放羊娃只能眼睁睁地看着，\n什么也做不了……",
        "body_size": 26,
    },
    {
        "draw": draw_scene_boy_sad,
        "title": "",
        "body": "放羊娃后悔极了。😢\n他再也不敢说谎了。\n可是羊已经回不来了……",
        "body_size": 26,
    },
    {
        "draw": draw_scene_lesson,
        "title": "小朋友们，记住哦！",
        "body": "🌟 说谎话是不对的！\n🌟 别人会不再相信你。\n🌟 做一个诚实的好孩子！❤️",
        "title_size": 32,
        "body_size": 26,
    },
    {
        "draw": draw_ending,
        "title": "— 故事讲完啦 —",
        "body": "谢谢小朋友们！\n你们要做诚实的好孩子哦！\n🌈✨👏",
        "title_size": 36,
        "body_size": 26,
    },
]


# ---------------------------------------------------------------------------
# PPT 生成
# ---------------------------------------------------------------------------

def _make_image(draw_fn) -> BytesIO:
    """生成一张 960×540 插图，返回内存中的 PNG。"""
    img = Image.new("RGB", (IMG_W, IMG_H), WHITE)
    draw_fn(img)
    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def build_ppt():
    prs = Presentation()
    # 16:9 尺寸
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    for slide_data in SLIDES:
        layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(layout)

        # 背景图
        img_buf = _make_image(slide_data["draw"])
        pic = slide.shapes.add_picture(img_buf, Emu(0), Emu(0), slide_w, slide_h)

        # 标题文本 (上方)
        title_text = slide_data.get("title", "")
        if title_text:
            title_size = slide_data.get("title_size", 36)
            txBox = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3), slide_w - Inches(1), Inches(1.2)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(title_size)
            p.font.bold = True
            p.font.color.rgb = RGBColor(180, 30, 30)
            p.alignment = PP_ALIGN.CENTER

        # 正文文本 (下方半透明框)
        body_text = slide_data.get("body", "")
        if body_text:
            body_size = slide_data.get("body_size", 22)
            # 半透明白色矩形背景可以通过形状实现
            box_top = slide_h - Inches(2.8)
            box_left = Inches(0.3)
            box_width = slide_w - Inches(0.6)
            box_height = Inches(2.5)

            # 白色底框
            shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE.RECTANGLE
                box_left, box_top, box_width, box_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            # 设置 75% 不透明度
            a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
            spPr = shape._element.find(f"{{{a_ns}}}spPr")
            if spPr is not None:
                solidFill = spPr.find(f"{{{a_ns}}}solidFill")
                if solidFill is not None:
                    srgbClr = solidFill.find(f"{{{a_ns}}}srgbClr")
                    if srgbClr is not None:
                        alpha = etree.SubElement(srgbClr, f"{{{a_ns}}}alpha")
                        alpha.set("val", "75000")
            shape.line.fill.background()

            # 文字
            tf = shape.text_frame
            tf.word_wrap = True
            for i, line in enumerate(body_text.split("\n")):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(body_size)
                p.font.bold = True
                p.font.color.rgb = RGBColor(50, 50, 80)
                p.alignment = PP_ALIGN.CENTER

    prs.save(OUTPUT_FILE)
    print(f"✅ PPT 已生成: {OUTPUT_FILE}  ({len(SLIDES)} 页)")


if __name__ == "__main__":
    build_ppt()
